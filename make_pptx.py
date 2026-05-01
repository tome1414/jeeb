from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color Palette ──────────────────────────────────────────────
NAVY    = RGBColor(0x0A, 0x0A, 0x2E)
BLUE    = RGBColor(0x00, 0x66, 0xFF)
BLUE_DIM= RGBColor(0xE0, 0xEB, 0xFF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF5, 0xF5, 0xF3)
GRAY    = RGBColor(0x5A, 0x5A, 0x56)
DARK    = RGBColor(0x0A, 0x0A, 0x0A)
GREEN   = RGBColor(0x00, 0xBB, 0x77)
AMBER   = RGBColor(0xFF, 0x99, 0x00)

W = Inches(10)   # slide width
H = Inches(5.625)  # slide height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))  # MSO_SHAPE_TYPE.RECTANGLE=1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=16, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def multiline_txt(slide, lines, x, y, w, h, size=14, color=DARK,
                  align=PP_ALIGN.LEFT, line_spacing_pt=6):
    """lines: list of (text, bold, size_override) tuples"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, sz = item, False, size
        elif len(item) == 2:
            text, bold = item; sz = size
        else:
            text, bold, sz = item

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        # spacing
        p.space_after = Pt(line_spacing_pt)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox

def card(slide, x, y, w, h, header_color=BLUE, bg_color=WHITE,
         header_text="", header_size=11, body_lines=None, body_size=11):
    """Draw a card with a colored top accent bar"""
    # background
    rect(slide, x, y, w, h, fill_color=bg_color, line_color=RGBColor(0xE0,0xE0,0xDC), line_width=Pt(0.75))
    # top accent bar
    rect(slide, x, y, w, 0.055, fill_color=header_color)
    # header text
    if header_text:
        txt(slide, header_text, x+0.12, y+0.07, w-0.2, 0.28,
            size=header_size, bold=True, color=header_color, align=PP_ALIGN.LEFT)
    # body
    if body_lines:
        by = y + 0.35
        for line in body_lines:
            if isinstance(line, tuple):
                t, sz = line
            else:
                t, sz = line, body_size
            txt(slide, t, x+0.12, by, w-0.24, 0.25, size=sz, color=GRAY)
            by += 0.22

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════
s1 = blank_slide(prs)
bg(s1, NAVY)

# diagonal accent shape (top-right corner flair)
accent = s1.shapes.add_shape(1, Inches(7.2), Inches(0), Inches(2.8), Inches(2.2))
accent.fill.solid(); accent.fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x99)
accent.line.fill.background()

# Bottom bar
rect(s1, 0, 4.9, 10, 0.725, fill_color=BLUE)

# JEEB big logo text
txt(s1, "JEEB", 0.6, 0.9, 5, 1.4, size=96, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Subtitle
txt(s1, "タイ発・飲食店ディスカバリープラットフォーム",
    0.6, 2.4, 8, 0.55, size=20, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.LEFT)

# Sub-subtitle
txt(s1, "エンジニア向け開発ブリーフィング資料",
    0.6, 3.05, 8, 0.45, size=14, color=RGBColor(0x77, 0x99, 0xCC), align=PP_ALIGN.LEFT)

# Tagline in bottom bar
txt(s1, "Find your perfect table in Bangkok",
    0.6, 4.97, 7, 0.4, size=13, italic=True, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.LEFT)

# Version tag
txt(s1, "v1.0  2026", 8.5, 5.0, 1.3, 0.35, size=11, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — 何を作るのか
# ══════════════════════════════════════════════════════════════════
s2 = blank_slide(prs)
bg(s2, LIGHT)

# Left blue accent bar
rect(s2, 0, 0, 0.08, 5.625, fill_color=BLUE)

# Title
txt(s2, "何を作るのか", 0.3, 0.25, 6, 0.65, size=30, bold=True, color=NAVY)
txt(s2, "プロジェクト概要", 0.3, 0.88, 5, 0.32, size=12, color=GRAY)

# Right callout box
rect(s2, 6.8, 0.8, 2.9, 2.2, fill_color=NAVY)
txt(s2, "\u300c", 6.98, 0.85, 0.4, 0.5, size=28, color=BLUE, bold=True)
txt(s2, "タイに食べログがない。", 6.95, 1.22, 2.6, 0.38, size=13, bold=True, color=WHITE)
txt(s2, "それが最大のチャンス。", 6.95, 1.55, 2.6, 0.38, size=13, bold=True, color=WHITE)
txt(s2, "\u300d", 9.05, 2.4, 0.5, 0.45, size=28, color=BLUE, bold=True, align=PP_ALIGN.RIGHT)

# 4 bullet rows
bullets = [
    ("🗺️", "多言語飲食店ポータル", "バンコク起点 → タイ全土 → アジアへ。全業態対応。"),
    ("🤖", "AIコンシェルジュ検索", "会話形式で最適な店を提案（Claude API連携）"),
    ("🏪", "店舗向け SaaS", "月額 2,000〜5,000 THB で多言語ページ + MEO + 予約"),
    ("📱", "Web → PWA → ネイティブ", "段階的にアプリ化。APIは最初からモバイル対応設計"),
]
for i, (icon, title, desc) in enumerate(bullets):
    by = 1.3 + i * 1.0
    # blue icon circle bg
    rect(s2, 0.25, by, 0.55, 0.52, fill_color=BLUE_DIM)
    txt(s2, icon, 0.27, by+0.02, 0.5, 0.45, size=20, align=PP_ALIGN.CENTER)
    txt(s2, title, 0.92, by, 3.5, 0.3, size=14, bold=True, color=NAVY)
    txt(s2, desc,  0.92, by+0.3, 5.6, 0.28, size=11, color=GRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — 既にここまで動いている
# ══════════════════════════════════════════════════════════════════
s3 = blank_slide(prs)
bg(s3, LIGHT)
rect(s3, 0, 0, 0.08, 5.625, fill_color=BLUE)

txt(s3, "既にここまで動いている", 0.3, 0.22, 8, 0.6, size=28, bold=True, color=NAVY)
txt(s3, "現行モックアップ構成（GitHub管理）", 0.3, 0.8, 7, 0.3, size=12, color=GRAY)

cards_data = [
    ("🏠  index.html",      "トップページ",    ["カード一覧・地図表示", "会員システム連携", "フィルター・お気に入り"]),
    ("🤖  search.html",     "AIコンシェルジュ",["6シナリオカード", "音声入力対応", "会話→結果画面"]),
    ("🍽️  restaurant.html", "店舗詳細ページ",  ["比較モーダル", "前払い決済UI", "LINE問い合わせ"]),
    ("👤  profile.html",    "会員プロフィール", ["Silver/Gold/Blackランク", "お気に入り・訪問履歴", "会員特典表示"]),
]
positions = [(0.22, 1.18), (5.12, 1.18), (0.22, 3.12), (5.12, 3.12)]
for (hdr, sub, items), (cx, cy) in zip(cards_data, positions):
    rect(s3, cx, cy, 4.68, 1.72, fill_color=WHITE,
         line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
    rect(s3, cx, cy, 4.68, 0.07, fill_color=BLUE)
    txt(s3, hdr, cx+0.15, cy+0.1, 4.3, 0.32, size=13, bold=True, color=NAVY)
    txt(s3, sub, cx+0.15, cy+0.42, 4.3, 0.28, size=11, color=BLUE, bold=True)
    for j, item in enumerate(items):
        txt(s3, "• " + item, cx+0.15, cy+0.7+j*0.27, 4.3, 0.25, size=10, color=GRAY)

# Bottom note
rect(s3, 0.22, 4.93, 9.56, 0.42, fill_color=BLUE_DIM)
txt(s3, "＋ 既存アセット: soi24-menu（5言語対応デジタルメニュー）— 統合・拡張予定",
    0.4, 4.96, 9.2, 0.36, size=11, color=NAVY, bold=True)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — マネタイズ構造
# ══════════════════════════════════════════════════════════════════
s4 = blank_slide(prs)
bg(s4, LIGHT)
rect(s4, 0, 0, 0.08, 5.625, fill_color=BLUE)

txt(s4, "マネタイズ構造", 0.3, 0.22, 6, 0.6, size=28, bold=True, color=NAVY)
txt(s4, "B2B（店舗側）+ B2C（ユーザー側）の2軸収益", 0.3, 0.8, 7, 0.3, size=12, color=GRAY)

# B2B column
rect(s4, 0.22, 1.18, 4.55, 3.5, fill_color=NAVY)
txt(s4, "店舗側  B2B", 0.42, 1.3, 4.0, 0.42, size=15, bold=True, color=WHITE)
rect(s4, 0.22, 1.68, 4.55, 0.03, fill_color=BLUE)

b2b = [
    ("💳  ¥2,000 THB / 月", "多言語ページ + MEO対応"),
    ("💳  ¥5,000 THB / 月", "上記 + 予約システム連携"),
    ("⚡  前払い決済",        "3日入金（競合は2〜4週間）"),
    ("📈  上位表示オプション","＋3,000〜10,000 THB/月"),
]
for i, (t, d) in enumerate(b2b):
    by = 1.8 + i * 0.72
    txt(s4, t, 0.42, by,      4.1, 0.3, size=12, bold=True, color=WHITE)
    txt(s4, d, 0.42, by+0.28, 4.1, 0.25, size=10, color=RGBColor(0xAA, 0xCC, 0xFF))

# B2C column
rect(s4, 5.0, 1.18, 4.78, 3.5, fill_color=WHITE,
     line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.75))
rect(s4, 5.0, 1.18, 4.78, 0.07, fill_color=BLUE)
txt(s4, "ユーザー側  B2C", 5.2, 1.3, 4.3, 0.42, size=15, bold=True, color=NAVY)
rect(s4, 5.0, 1.68, 4.78, 0.03, fill_color=RGBColor(0xCC,0xCC,0xFF))

b2c = [
    ("🥈  Silver（無料）",   "検索・AIコンシェルジュ全機能"),
    ("🥇  Gold（〜250 THB/月）","加盟店3%割引・限定情報"),
    ("⬛  Black（招待制）",   "限定予約・専属コンシェルジュ"),
]
for i, (t, d) in enumerate(b2c):
    by = 1.8 + i * 0.82
    txt(s4, t, 5.2, by,      4.4, 0.3, size=12, bold=True, color=NAVY)
    txt(s4, d, 5.2, by+0.3,  4.4, 0.25, size=10, color=GRAY)

# Bottom highlight
rect(s4, 0.22, 4.78, 9.56, 0.55, fill_color=BLUE)
txt(s4, "🎯  加盟 200 店で月次黒字化の見込み　｜　LTV / CAC 比率：推定 8〜15 倍",
    0.42, 4.84, 9.1, 0.38, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — 技術スタック
# ══════════════════════════════════════════════════════════════════
s5 = blank_slide(prs)
bg(s5, LIGHT)
rect(s5, 0, 0, 0.08, 5.625, fill_color=BLUE)

txt(s5, "推奨技術スタック", 0.3, 0.22, 7, 0.6, size=28, bold=True, color=NAVY)
txt(s5, "スケールを見据えたモダン構成", 0.3, 0.8, 7, 0.3, size=12, color=GRAY)

cols = [
    ("🖥️  Frontend",   ["Next.js 14 (App Router)", "Tailwind CSS", "Leaflet.js / Google Maps", "Web Speech API（音声入力）"]),
    ("⚙️  Backend & AI",["Node.js + Fastify", "PostgreSQL + Redis", "Elasticsearch（店舗検索）", "Anthropic Claude API"]),
    ("☁️  Infrastructure",["AWS / GCP", "Omise 決済（タイ特化）", "LINE Messaging API", "Firebase（PWA・Push通知）"]),
]
col_x = [0.22, 3.55, 6.88]
for (hdr, items), cx in zip(cols, col_x):
    rect(s5, cx, 1.18, 3.12, 3.9, fill_color=WHITE,
         line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.5))
    rect(s5, cx, 1.18, 3.12, 0.5, fill_color=NAVY)
    txt(s5, hdr, cx+0.15, 1.25, 2.8, 0.35, size=13, bold=True, color=WHITE)
    for i, item in enumerate(items):
        iy = 1.82 + i * 0.72
        rect(s5, cx+0.15, iy, 2.82, 0.5, fill_color=LIGHT,
             line_color=RGBColor(0xE0,0xE0,0xDC), line_width=Pt(0.5))
        txt(s5, item, cx+0.28, iy+0.1, 2.6, 0.3, size=11, color=DARK)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — メニュー × リアクション
# ══════════════════════════════════════════════════════════════════
s6 = blank_slide(prs)
bg(s6, LIGHT)
rect(s6, 0, 0, 0.08, 5.625, fill_color=BLUE)

txt(s6, "メニュー × ソーシャルリアクション", 0.3, 0.22, 9, 0.6, size=26, bold=True, color=NAVY)
txt(s6, "soi24-menu（5言語対応）を統合し、料理1品ごとにSNS的インタラクションを付加",
    0.3, 0.8, 9.4, 0.3, size=12, color=GRAY)

feat_cards = [
    ("❤️", "いいね & 食べたい",   ["メニュー1品ごとにリアクション", "マイウィッシュリストに蓄積", "食べた後→レビューへ自動誘導"]),
    ("🌟", "著名人バッジ",         ["インフルエンサー公式連携", "「○○さんが食べたい」表示", "Verified済みで信頼性担保"]),
    ("📸", "SNS話題バッジ",        ["TikTok・Instagram投稿と紐づけ", "コミュニティ報告→一定数で自動付与", "「バズってる一品」として訴求"]),
    ("📍", "来店誘導プッシュ通知",  ["食べたい品の店舗に接近を検知", "「今すぐ食べられます」通知", "GPS連動（PWA・ネイティブアプリ）"]),
]
positions6 = [(0.22, 1.18), (5.12, 1.18), (0.22, 3.2), (5.12, 3.2)]
for (icon, title, desc_list), (cx, cy) in zip(feat_cards, positions6):
    rect(s6, cx, cy, 4.68, 1.85, fill_color=WHITE,
         line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.5))
    rect(s6, cx, cy, 4.68, 0.07, fill_color=BLUE)
    txt(s6, icon + "  " + title, cx+0.15, cy+0.1, 4.3, 0.35, size=13, bold=True, color=NAVY)
    for j, d in enumerate(desc_list):
        txt(s6, "• " + d, cx+0.15, cy+0.5+j*0.35, 4.3, 0.32, size=10, color=GRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — ロードマップ
# ══════════════════════════════════════════════════════════════════
s7 = blank_slide(prs)
bg(s7, LIGHT)
rect(s7, 0, 0, 0.08, 5.625, fill_color=BLUE)

txt(s7, "4フェーズ 開発ロードマップ", 0.3, 0.22, 8, 0.6, size=28, bold=True, color=NAVY)
txt(s7, "α版でまず事例を作り、数字で拡大を正当化する", 0.3, 0.8, 7, 0.3, size=12, color=GRAY)

# Timeline bar
rect(s7, 0.5, 2.3, 9.0, 0.12, fill_color=RGBColor(0xCC,0xCC,0xCC))

phases = [
    ("α版", "〜2ヶ月", "店舗掲載・検索\n会員登録", "50店舗登録", BLUE),
    ("β版", "〜4ヶ月", "AI検索・予約\nLINE連携", "MAU 1,000", RGBColor(0x00,0x88,0xDD)),
    ("v1.0", "〜6ヶ月", "決済・多言語\n管理画面", "加盟100店・黒字化", GREEN),
    ("v2.0", "〜12ヶ月","バンコク全域\n分析レポート", "加盟500店", AMBER),
]
phase_x = [0.55, 2.88, 5.2, 7.52]
for (label, period, scope, kpi, color), px in zip(phases, phase_x):
    # top box
    rect(s7, px, 1.12, 2.0, 1.05, fill_color=color)
    txt(s7, label,  px+0.1, 1.18, 1.8, 0.38, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s7, period, px+0.1, 1.52, 1.8, 0.28, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    # dot on timeline
    dot = s7.shapes.add_shape(9, Inches(px+0.85), Inches(2.24), Inches(0.3), Inches(0.3))  # OVAL=9
    dot.fill.solid(); dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    # bottom box
    rect(s7, px, 2.55, 2.0, 1.05, fill_color=WHITE,
         line_color=color, line_width=Pt(1.5))
    scope_lines = scope.split('\n')
    for li, sl in enumerate(scope_lines):
        txt(s7, sl, px+0.1, 2.6+li*0.28, 1.8, 0.27, size=10, color=DARK, align=PP_ALIGN.CENTER)
    # arrow between phases
    if px < 7.5:
        txt(s7, "▶", px+2.05, 2.18, 0.28, 0.28, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    # KPI badge
    rect(s7, px, 3.72, 2.0, 0.4, fill_color=BLUE_DIM)
    txt(s7, "🎯 " + kpi, px+0.05, 3.76, 1.9, 0.3, size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — Stretch Goals
# ══════════════════════════════════════════════════════════════════
s8 = blank_slide(prs)
bg(s8, NAVY)

txt(s8, "将来の差別化機能", 0.4, 0.22, 8, 0.6, size=28, bold=True, color=WHITE)
txt(s8, "実現すれば圧倒的な優位性 — 今は不確実、でも狙う価値がある",
    0.4, 0.8, 9, 0.3, size=12, color=RGBColor(0xAA, 0xCC, 0xFF), italic=True)

stretch = [
    ("🔮", "AI メニュー自動デジタル化", "紙メニュー撮影 → 即多言語化 (Claude Vision)"),
    ("👥", "ソーシャルダイニング",       "「今夜一緒に食べませんか」マッチング機能"),
    ("⏱️", "リアルタイム空席通知",       "キャンセル枠の自動再販売・プッシュ通知"),
    ("🤳", "料理 AI カメラ判定",         "写真から食べられる店を逆検索"),
    ("💳", "JEEB コイン",               "来店・レビューでポイント獲得→加盟店で使用"),
    ("📊", "店舗 AI アドバイス",         "「火曜18時に空席多→クーポン推奨」自動提案"),
]
positions8 = [
    (0.22, 1.18), (5.12, 1.18),
    (0.22, 2.75), (5.12, 2.75),
    (0.22, 4.32), (5.12, 4.32),
]
for (icon, title, desc), (cx, cy) in zip(stretch, positions8):
    rect(s8, cx, cy, 4.68, 1.35, fill_color=RGBColor(0x12,0x12,0x44),
         line_color=RGBColor(0x00,0x44,0xAA), line_width=Pt(0.75))
    txt(s8, icon, cx+0.15, cy+0.15, 0.6, 0.5, size=20, align=PP_ALIGN.CENTER)
    txt(s8, title, cx+0.8, cy+0.12, 3.7, 0.35, size=12, bold=True, color=WHITE)
    txt(s8, desc,  cx+0.8, cy+0.47, 3.7, 0.5,  size=10, color=RGBColor(0xAA,0xCC,0xFF))

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — Web → App Evolution
# ══════════════════════════════════════════════════════════════════
s9 = blank_slide(prs)
bg(s9, LIGHT)
rect(s9, 0, 0, 0.08, 5.625, fill_color=BLUE)

txt(s9, "Web から アプリ への Evolution", 0.3, 0.22, 8, 0.6, size=26, bold=True, color=NAVY)
txt(s9, "段階的移行でリスクを最小化し、ユーザー体験を継続的に向上", 0.3, 0.8, 9, 0.3, size=12, color=GRAY)

stages = [
    ("🌐", "Web App", "（現在）",   ["モバイルファースト設計", "GitHub管理・即日デプロイ", "ブラウザで全機能利用可"]),
    ("📲", "PWA",      "（β版後）",  ["ホーム画面に追加可能", "オフライン対応（メニューキャッシュ）", "プッシュ通知対応"]),
    ("📱", "Native App","（v2.0）",  ["React Native / Flutter", "GPS・カメラ・Apple Pay", "App Store 申請・公開"]),
]
stage_colors = [BLUE, RGBColor(0x00,0x99,0x88), GREEN]
stage_x = [0.22, 3.55, 6.88]

for (icon, title, period, items), color, sx in zip(stages, stage_colors, stage_x):
    rect(s9, sx, 1.18, 3.12, 3.7, fill_color=WHITE,
         line_color=color, line_width=Pt(1.5))
    rect(s9, sx, 1.18, 3.12, 0.85, fill_color=color)
    txt(s9, icon,   sx+0.15, 1.22, 2.8, 0.42, size=22, align=PP_ALIGN.CENTER, color=WHITE)
    txt(s9, title,  sx+0.1,  1.58, 2.9, 0.32, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s9, period, sx+0.1,  1.87, 2.9, 0.22, size=10, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
    for i, item in enumerate(items):
        txt(s9, "✓  " + item, sx+0.2, 2.18+i*0.5, 2.75, 0.42, size=11, color=DARK)

    # arrow
    if sx < 6.8:
        txt(s9, "▶", sx+3.17, 2.55, 0.35, 0.35, size=18, color=BLUE, align=PP_ALIGN.CENTER)

# Bottom important note
rect(s9, 0.22, 5.05, 9.56, 0.42, fill_color=BLUE_DIM)
txt(s9, "⚠️  バックエンド API は最初からモバイル対応の RESTful 設計で構築すること",
    0.42, 5.08, 9.1, 0.35, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — チーム・コスト概算
# ══════════════════════════════════════════════════════════════════
s10 = blank_slide(prs)
bg(s10, NAVY)

# accent top strip
rect(s10, 0, 0, 10, 0.12, fill_color=BLUE)

txt(s10, "開発リソース概算", 0.4, 0.28, 7, 0.6, size=28, bold=True, color=WHITE)
txt(s10, "α版（2ヶ月）を起点に、事例と数字を積み上げる", 0.4, 0.86, 8, 0.3, size=12,
    color=RGBColor(0xAA,0xCC,0xFF), italic=True)

# Table header
headers = ["ロール", "人数", "月単価目安（THB）"]
col_ws  = [3.8, 1.5, 3.0]
col_xs  = [0.3, 4.2, 5.8]
hy = 1.28
for hdr, cx, cw in zip(headers, col_xs, col_ws):
    rect(s10, cx, hy, cw, 0.45, fill_color=BLUE)
    txt(s10, hdr, cx+0.1, hy+0.07, cw-0.15, 0.3, size=12, bold=True, color=WHITE)

rows = [
    ("フルスタックエンジニア", "2〜3名", "15〜30 万"),
    ("UIデザイナー",           "1名",    " 8〜15 万"),
    ("QAエンジニア",           "1名",    " 6〜10 万"),
    ("合計（α版 2ヶ月）",     "—",      "約 60〜150 万"),
]
row_colors = [RGBColor(0x12,0x12,0x44), RGBColor(0x10,0x10,0x38),
              RGBColor(0x12,0x12,0x44), RGBColor(0x00,0x33,0x88)]
text_colors= [WHITE, WHITE, WHITE, RGBColor(0xFF, 0xFF, 0xAA)]
for i, (role, count, cost) in enumerate(rows):
    ry = 1.73 + i * 0.5
    for cx, cw in zip(col_xs, col_ws):
        rect(s10, cx, ry, cw, 0.48, fill_color=row_colors[i],
             line_color=RGBColor(0x22,0x22,0x66), line_width=Pt(0.5))
    txt(s10, role,  col_xs[0]+0.12, ry+0.1, col_ws[0]-0.2, 0.3, size=12, color=text_colors[i], bold=(i==3))
    txt(s10, count, col_xs[1]+0.1,  ry+0.1, col_ws[1]-0.15,0.3, size=12, color=text_colors[i], bold=(i==3), align=PP_ALIGN.CENTER)
    txt(s10, cost,  col_xs[2]+0.1,  ry+0.1, col_ws[2]-0.15,0.3, size=12, color=text_colors[i], bold=(i==3), align=PP_ALIGN.RIGHT)

# Closing message
rect(s10, 0.3, 4.15, 9.4, 1.2, fill_color=RGBColor(0x00,0x22,0x66))
rect(s10, 0.3, 4.15, 0.08, 1.2, fill_color=BLUE)
txt(s10, "まずは α 版で Prompong / Thonglor 50 店を獲得し、",
    0.55, 4.25, 9.0, 0.4, size=14, bold=True, color=WHITE)
txt(s10, "「継続率・予約数・オーナー満足度」の数字を揃える。それが全ての起点。",
    0.55, 4.62, 9.0, 0.4, size=13, color=RGBColor(0xAA,0xCC,0xFF))
txt(s10, "JEEB  2026", 7.5, 5.2, 2.2, 0.3, size=11, color=RGBColor(0x55,0x77,0xAA),
    align=PP_ALIGN.RIGHT, italic=True)

# ══════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════
out = "/Users/saotome/Claude/JEEB_Portal/JEEB_Engineer_Brief.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
